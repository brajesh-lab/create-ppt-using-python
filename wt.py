from wand.image import Image
import os
from pptx.text.text import Font
from pptx import Presentation
from pptx.util import Cm,Inches,Pt

images=[] 
def getFiles(path): 
    for file in os.listdir(path): 
        if file.endswith(".jpg"): 
            images.append(os.path.join(path, file))
path="D:\\ppt"
getFiles(path)
i=0
# Import the image
def water(image):
 i=0
 for imgee in images[1:]:

    with Image(filename =imgee) as image:
     
    
        # Import the watermark image
     with Image(filename ='D:\\nike_black.png') as water:
        
        water.resize(1500,600)
        with image.clone() as watermark:
            # left as 10 and top as 20
            watermark.watermark(water, 0.1, 20, 20)
                # Save the image
            watermark.save(filename="D:\\water\Image%d.jpg"%i)
            i=i+1
            
water(images)
#ppt ..................................................................................................................


def create_presentation(filename, content):
    "Create a presentation using a multiline string"
    # transform the multiline in a list
    path = ["D:\\water\Image0.jpg","D:\\water\Image1.jpg","D:\\water\Image2.jpg","D:\\water\Image3.jpg","D:\\water\Image4.jpg"]
    content2 = []
    i=0
    #split text into title and subtitle
    for line in content.split("\n\n"):
        line = line.split("\n")
        content2.append(line)
    prs = Presentation()
    layout = prs.slide_layouts[8] 

    for lst in content2:
        slide = prs.slides #new slide
        slide = slide.add_slide(layout)
        picture = slide.shapes.placeholders[1]  #picture place holder
        title = slide.shapes.placeholders[0]   
        title.width=Cm(15)
        title.height=Cm(1)
        title.top = Cm(1)
        title.left=Cm(10)
        title.text = lst[0]
        
       # picture = slide.placeholders[1]
        pic = picture.insert_picture(path[i])
        pic.top=Cm(5)
        pic.crop_top = 0 
        pic.crop_left = 0 
        pic.crop_bottom = 0
        pic.crop_right = 0
        pic.width= Inches(6)
        pic.height= Inches(6)
        pic.left=Cm(3)
        i+=1
        subtitle = slide.placeholders[2]
        subtitle.width=Cm(15)
        subtitle.height=Cm(1)
        subtitle.top = Cm(2)
        subtitle.left=Cm(3)
        subtitle.text = lst[1]

    # save and launch the file
    prs.save(filename)
    os.startfile(filename)


content = """BREAKFAST 
VEG. 

WORKING 
Important

GUEST ROOM
with green env.

CAMERA
With high Quality

DECORATION
ROOM DECORATION"""



create_presentation("D:\\example.pptx", content)